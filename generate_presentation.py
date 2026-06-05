#!/usr/bin/env python3
"""
generate_presentation.py

Generates a professional 28-slide PowerPoint (.pptx) presentation
about the HazardMap real-time earthquake hazard mapping application.

Run:
    source .pptx_venv/bin/activate
    python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

# ── Branding Colors ──────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x14, 0x1E)  # Deep navy
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)  # Cyan accent
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)  # Green
ACCENT_ORANGE= RGBColor(0xF9, 0x73, 0x16)  # Orange
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)  # Red
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB8, 0xC4)
MED_GRAY     = RGBColor(0x70, 0x78, 0x84)
CARD_BG      = RGBColor(0x15, 0x1F, 0x2E)  # Card background
SLIDE_BG     = RGBColor(0x0B, 0x11, 0x1A)  # Slightly darker slide bg

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helper Functions ─────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set solid fill background for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=Inches(0.6), top=None, width=Inches(1.5), height=Inches(0.06), color=ACCENT_CYAN):
    """Add a thin colored accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_text(slide, text, left=Inches(0.6), top=Inches(0.4), width=Inches(12), height=Inches(0.8),
                   font_size=36, color=WHITE, bold=True):
    """Add a title text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    return txBox


def add_subtitle_text(slide, text, left=Inches(0.6), top=Inches(1.1), width=Inches(12), height=Inches(0.6),
                      font_size=18, color=LIGHT_GRAY):
    """Add subtitle text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return txBox


def add_body_text(slide, text, left=Inches(0.6), top=Inches(2.0), width=Inches(12), height=Inches(5),
                  font_size=14, color=LIGHT_GRAY, line_spacing=1.3):
    """Add body text with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
        if line.strip().startswith("•"):
            p.level = 0
        elif line.strip().startswith("–") or line.strip().startswith("-"):
            p.level = 1
    return txBox


def add_card(slide, left, top, width, height, title, body_lines, accent_color=ACCENT_CYAN):
    """Add a card-style content box."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(0x25, 0x35, 0x4A)
    card.line.width = Pt(1)

    # Card title
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = accent_color
    p.font.bold = True
    p.font.name = "Calibri"

    # Card body
    body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7))
    tf2 = body_box.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = "Calibri"
        p2.space_after = Pt(2)


def add_table(slide, left, top, width, height, headers, rows, accent_color=ACCENT_CYAN):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Style header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3F)

    # Style rows
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if i % 2 == 0 else RGBColor(0x12, 0x1C, 0x2B)

    return table_shape


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_section_number(slide, number, top=Inches(0.35)):
    """Add a small slide number badge."""
    txBox = slide.shapes.add_textbox(Inches(12.2), top, Inches(0.8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number:02d}"
    p.font.size = Pt(14)
    p.font.color.rgb = MED_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


# ── Presentation Builder ─────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ================================================================
    # SLIDE 1: Title Slide
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, DARK_BG)

    # Large title
    add_title_text(slide, "HazardMap", left=Inches(0.8), top=Inches(1.8), font_size=54, color=ACCENT_CYAN)
    add_subtitle_text(slide, "Real-Time Earthquake Hazard Mapping & Simulation Engine",
                      left=Inches(0.8), top=Inches(2.8), font_size=24, color=WHITE)
    add_accent_bar(slide, left=Inches(0.8), top=Inches(3.5), width=Inches(3), color=ACCENT_CYAN)

    add_body_text(slide, 
        "A scientific-grade web application for real-time seismic hazard assessment\n"
        "covering the Indian subcontinent and surrounding regions.\n\n"
        "Version 3.0.0  •  Built with FastAPI, React, PostgreSQL, MapLibre GL",
        left=Inches(0.8), top=Inches(3.8), font_size=16, color=LIGHT_GRAY)

    add_speaker_notes(slide,
        "Welcome to the HazardMap presentation. This application provides real-time earthquake "
        "hazard mapping and simulation capabilities for the Indian subcontinent. "
        "It combines live seismic data feeds, ground motion prediction equations, "
        "soil amplification modeling, and interactive web-based visualization.")

    # ================================================================
    # SLIDE 2: Executive Summary
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 2)
    add_title_text(slide, "Executive Summary")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(2.4),
        "What It Does",
        ["• Real-time earthquake monitoring from",
         "  USGS and NCS (India) feeds",
         "• Computes PGA (Peak Ground Acceleration)",
         "  for 6,500+ grid cells across India",
         "• Generates contour maps and impact",
         "  summaries within seconds"])

    add_card(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(2.4),
        "Key Objectives",
        ["• Provide sub-minute hazard assessment",
         "  after earthquake detection",
         "• Integrate Vs30 soil amplification for",
         "  site-specific ground motion estimates",
         "• Enable manual \"what-if\" simulations",
         "  for disaster preparedness planning"])

    add_card(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(2.4),
        "Main Capabilities",
        ["• Live event ingestion & deduplication",
         "• Region-specific GMPE model selection",
         "• NEHRP site classification from raster data",
         "• Vectorized PGA computation (NumPy)",
         "• WebSocket-based real-time UI updates",
         "• CSV, JSON, GeoJSON export"])

    add_card(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.6),
        "System Highlights",
        ["• Covers 3 tectonic regions: Himalayan, Northeast India, Peninsular — each with tuned GMPE coefficients",
         "• Dual data sources: USGS global feed + NCS India (scraped) — with spatial-temporal deduplication",
         "• Full-stack Docker deployment: React frontend + FastAPI backend + PostgreSQL database",
         "• End-to-end latency: earthquake detected → hazard map rendered on client in < 10 seconds",
         "• 6,500+ nationwide 20km grid cells, each with Vs30, site class, soil factor, base PGA, and amplified PGA"])

    add_speaker_notes(slide,
        "HazardMap is a full-stack real-time seismic hazard assessment platform. "
        "It ingests earthquake data from two sources (USGS and NCS India), runs ground motion "
        "prediction equations with regional calibration, applies Vs30-based soil amplification, "
        "generates smooth PGA contour maps, and pushes results to connected browser clients via WebSockets — "
        "all within seconds of an earthquake being detected.")

    # ================================================================
    # SLIDE 3: Table of Contents
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 3)
    add_title_text(slide, "Presentation Outline")
    add_accent_bar(slide, top=Inches(1.1))

    col1 = [
        "01. Title",
        "02. Executive Summary",
        "03. Outline",
        "04. Project Goals",
        "05. System Overview",
        "06. Application Architecture",
        "07. Technology Stack",
        "08. UI Walkthrough",
        "09. Feature Overview",
        "10. Interactive Map Features",
        "11. Earthquake Simulation Module",
        "12. Model Mathematics (GMPE)",
        "13. Model Mathematics (Soil & Distance)",
        "14. Scientific Foundations",
    ]
    col2 = [
        "15. Data Pipeline",
        "16. Data Sources",
        "17. External APIs & Services",
        "18. Database Design",
        "19. Core Algorithms",
        "20. Map Rendering Process",
        "21. User Workflow Example",
        "22. Performance & Scalability",
        "23. Security Considerations",
        "24. Challenges Encountered",
        "25. Limitations",
        "26. Key Learnings",
        "27. Conclusion",
        "28. References",
    ]

    add_body_text(slide, "\n".join(col1), left=Inches(0.8), top=Inches(1.6), width=Inches(5.5),
                  font_size=13, color=LIGHT_GRAY)
    add_body_text(slide, "\n".join(col2), left=Inches(6.8), top=Inches(1.6), width=Inches(5.5),
                  font_size=13, color=LIGHT_GRAY)

    add_speaker_notes(slide,
        "This presentation covers 28 slides spanning the full scope of the project — "
        "from system architecture and mathematics to deployment and future work.")

    # ================================================================
    # SLIDE 4: Project Goals
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 4)
    add_title_text(slide, "Project Goals")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0),
        "Functional Goals",
        ["• Ingest real-time earthquake data from USGS and NCS India",
         "• Automatically compute seismic hazard maps within seconds",
         "• Display interactive PGA heatmaps on a web-based map interface",
         "• Provide district-level and state-level impact summaries",
         "• Support manual earthquake simulations with custom parameters",
         "• Enable GeoJSON, CSV, and JSON data export",
         "• Push real-time updates to all connected clients via WebSocket",
         "• Classify tectonic regions and apply region-specific GMPE models",
         "• Integrate soil amplification using national Vs30 raster data",
         "• Deduplicate earthquake events across multiple data sources"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.0),
        "Technical Goals",
        ["• Vectorized NumPy computation for 6,500+ cells (no Python loops)",
         "• Asyncio-based concurrency: non-blocking event loop",
         "• Thread-safe PostgreSQL connection pool with psycopg2",
         "• Zero-downtime hot-reload in development via uvicorn StatReload",
         "• Docker Compose orchestration (frontend + backend + database)",
         "• Modular GMPE framework with abstract base class for extensibility",
         "• GeoTIFF raster caching with rasterio for Vs30 sampling",
         "• Matplotlib Agg backend for thread-safe contour generation",
         "• Spatial indexing using Shapely for boundary and region checks",
         "• Sub-10-second end-to-end pipeline from detection to display"],
        accent_color=ACCENT_CYAN)

    add_speaker_notes(slide,
        "The functional goals focus on real-time earthquake monitoring and hazard computation for India, "
        "while the technical goals ensure the system is performant, scalable, and scientifically grounded. "
        "Key technical decisions include vectorized NumPy computation instead of Python loops, "
        "asyncio for non-blocking I/O, and Docker Compose for reproducible deployment.")

    # ================================================================
    # SLIDE 5: System Overview
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 5)
    add_title_text(slide, "System Overview")
    add_accent_bar(slide, top=Inches(1.1))

    add_subtitle_text(slide, "High-Level Data Flow: From Earthquake Detection to Hazard Visualization",
                      top=Inches(1.3), font_size=14, color=MED_GRAY)

    # Flow diagram using text boxes and arrows
    steps = [
        ("USGS / NCS\nData Feed", ACCENT_CYAN),
        ("Poller &\nNormalizer", ACCENT_GREEN),
        ("Deduplication\n& Filter", ACCENT_ORANGE),
        ("Simulation\nWorker", ACCENT_RED),
        ("WebSocket\nBroadcast", ACCENT_CYAN),
        ("Browser\nMap UI", ACCENT_GREEN),
    ]
    box_w = Inches(1.7)
    box_h = Inches(1.0)
    start_x = Inches(0.6)
    y = Inches(2.5)
    gap = Inches(0.35)

    for i, (label, color) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = "Calibri"

        # Arrow between boxes
        if i < len(steps) - 1:
            arrow_x = x + box_w + Inches(0.05)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y + Inches(0.35), Inches(0.25), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MED_GRAY
            arrow.line.fill.background()

    # Core Modules section
    add_card(slide, Inches(0.6), Inches(4.0), Inches(3.8), Inches(3.0),
        "Core Modules",
        ["• ingest/ — USGS + NCS polling, normalization, dedup",
         "• layers/pga/ — GMPE framework, region detection, PGA engine",
         "• soil/ — Vs30 raster loading, site classification, amplification",
         "• services/ — contour generation, impact aggregation",
         "• jobs/ — async queue + simulation worker",
         "• api/ — REST endpoints + WebSocket"],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(4.7), Inches(4.0), Inches(3.8), Inches(3.0),
        "User Workflow",
        ["1. User opens HazardMap in a browser",
         "2. WebSocket connects → LIVE status indicator",
         "3. Earthquake detected by poller → event queued",
         "4. Worker runs simulation → PGA computed",
         "5. Results broadcast to all clients via WS",
         "6. Map auto-updates with contour heatmap",
         "   OR: User manually runs simulation via UI"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(8.8), Inches(4.0), Inches(3.8), Inches(3.0),
        "Key Design Decisions",
        ["• asyncio.to_thread() for blocking matplotlib",
         "  calls — keeps event loop free for WS/poller",
         "• Pre-loaded nationwide grid (6,500 cells)",
         "  — deep-copied per simulation, not re-read",
         "• Single national Vs30 raster instead of",
         "  per-state tiles — simpler, single I/O call",
         "• Tectonic boundaries from real plate GeoJSON"],
        accent_color=ACCENT_ORANGE)

    add_speaker_notes(slide,
        "The system follows a clean pipeline: external earthquake feeds are polled every 60 seconds, "
        "normalized into a canonical format, deduplicated using spatial-temporal fingerprints, "
        "filtered for India relevance, and queued for simulation. The simulation worker computes PGA "
        "for all grid cells, applies soil amplification, generates contour maps, and broadcasts results "
        "via WebSocket to all connected browser clients. The entire pipeline runs in under 10 seconds.")

    # ================================================================
    # SLIDE 6: Application Architecture
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 6)
    add_title_text(slide, "Application Architecture")
    add_accent_bar(slide, top=Inches(1.1))

    # Architecture layers
    layers_data = [
        ("PRESENTATION LAYER", "React 19 + MapLibre GL JS + Zustand State", ACCENT_CYAN, Inches(1.6)),
        ("API GATEWAY", "FastAPI REST + WebSocket (/api/ws/live)", ACCENT_GREEN, Inches(2.5)),
        ("BUSINESS LOGIC", "SimulationRunner → PGAEngine → SoilEngine → ContourGenerator → ImpactAggregator", ACCENT_ORANGE, Inches(3.4)),
        ("DATA INGESTION", "USGS Poller + NCS Scraper → Normalizer → Deduplicator → Filter → Queue", ACCENT_RED, Inches(4.3)),
        ("PERSISTENCE", "PostgreSQL 16 (simulations, events, dedup_cache) + GeoTIFF Rasters (Vs30)", ACCENT_CYAN, Inches(5.2)),
    ]

    for label, desc, color, top in layers_data:
        # Layer box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(12), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        # Label
        txBox = slide.shapes.add_textbox(Inches(0.9), top + Inches(0.08), Inches(3), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = "Calibri"

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(0.9), top + Inches(0.35), Inches(11.5), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = "Calibri"

    # Deployment box
    add_card(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.9),
        "Deployment: Docker Compose",
        ["3 containers: frontend (Vite dev server :5173) → backend (uvicorn :8000) → PostgreSQL 16 (:5432)  |  Volumes: /app/data (grids, rasters, GeoJSON)"],
        accent_color=MED_GRAY)

    add_speaker_notes(slide,
        "The architecture follows a clean layered pattern. The presentation layer is a React SPA "
        "using MapLibre GL for map rendering and Zustand for state management. The API layer uses FastAPI "
        "with both REST endpoints and a WebSocket for real-time push. The business logic layer contains "
        "the core simulation pipeline. Data ingestion runs as background asyncio tasks. "
        "Everything is orchestrated via Docker Compose with PostgreSQL for persistence.")

    # ================================================================
    # SLIDE 7: Technology Stack
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 7)
    add_title_text(slide, "Technology Stack")
    add_accent_bar(slide, top=Inches(1.1))

    headers = ["Category", "Technology", "Version", "Purpose"]
    rows = [
        ["Frontend Framework", "React", "19.x", "Component-based UI"],
        ["Map Rendering", "MapLibre GL JS", "4.x", "WebGL vector map rendering"],
        ["State Management", "Zustand", "5.x", "Lightweight global state store"],
        ["Build Tool", "Vite", "6.x", "Fast HMR dev server + bundler"],
        ["Backend Framework", "FastAPI", "0.115+", "Async REST API + WebSocket"],
        ["Runtime", "Python", "3.10+", "Backend language"],
        ["Database", "PostgreSQL", "16-alpine", "ACID-compliant persistence"],
        ["DB Driver", "psycopg2", "2.9+", "Thread-safe connection pool"],
        ["Geospatial", "Shapely + GeoPandas", "2.x / 1.x", "Geometry operations, boundary checks"],
        ["Raster I/O", "rasterio", "1.4+", "GeoTIFF Vs30 sampling"],
        ["Computation", "NumPy + SciPy", "1.26+ / 1.14+", "Vectorized PGA math, interpolation"],
        ["Visualization", "Matplotlib", "3.9+", "Contour polygon generation"],
        ["Contour Export", "geojsoncontour", "0.4+", "matplotlib contourf → GeoJSON"],
        ["HTTP Client", "aiohttp", "3.10+", "Async polling of USGS/NCS feeds"],
        ["Deployment", "Docker Compose", "3.8", "Multi-container orchestration"],
    ]

    add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5), headers, rows)

    add_speaker_notes(slide,
        "The technology stack was chosen for a combination of developer velocity, scientific computing capability, "
        "and real-time web performance. FastAPI provides native async support for polling and WebSocket. "
        "NumPy enables vectorized PGA computation over 6,500 grid cells without Python loops. "
        "MapLibre GL renders vector tiles with WebGL for smooth map interactions. "
        "Docker Compose ensures the entire stack is reproducible across environments.")

    # ================================================================
    # SLIDE 8: UI Walkthrough
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 8)
    add_title_text(slide, "User Interface Walkthrough")
    add_accent_bar(slide, top=Inches(1.1))

    ui_elements = [
        ("Navbar", "Top bar with app branding (HazardMap v1.0.0), theme toggle (dark/light), notification bell, user avatar", ACCENT_CYAN),
        ("Sidebar", "Left icon rail — sections: Simulation Panel, Layers Panel, Alerts Panel, Disasters Panel", ACCENT_GREEN),
        ("Map View", "Full-screen MapLibre GL map with India boundary overlay, state borders, and grid cells. Supports satellite/terrain toggle.", ACCENT_ORANGE),
        ("Simulation Panel", "Click-to-place epicenter → set magnitude (1-10) + depth (1-100km) → Run Simulation button. Shows region auto-detection.", ACCENT_RED),
        ("Live Events Panel", "Real-time feed of earthquake events from USGS and NCS. Shows LIVE/OFFLINE indicator. Auto-updates via WebSocket.", ACCENT_CYAN),
        ("Map Legend", "Color-coded PGA intensity bands: No Effect → Light → Moderate → Strong → Very Strong → Severe → Violent", ACCENT_GREEN),
        ("Layers Panel", "Toggle: satellite imagery, terrain, India boundary, state boundaries, soil amplification overlay. Opacity sliders.", ACCENT_ORANGE),
    ]

    y_start = Inches(1.5)
    for i, (name, desc, color) in enumerate(ui_elements):
        y = y_start + i * Inches(0.8)
        # Colored dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.12), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # Name
        txBox = slide.shapes.add_textbox(Inches(0.95), y, Inches(2.2), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = "Calibri"

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(3.2), y, Inches(9.5), Inches(0.7))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = "Calibri"

    add_speaker_notes(slide,
        "The UI follows a dark-themed GIS dashboard layout. The full-screen map occupies the center, "
        "with a collapsible sidebar on the left for simulation controls, layer toggles, and live event feeds. "
        "The navbar provides branding and theme controls. The map legend shows PGA intensity color coding "
        "matching USGS ShakeMap standards. All UI updates are real-time via WebSocket push.")

    # ================================================================
    # SLIDE 9: Feature Overview
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 9)
    add_title_text(slide, "Feature Overview")
    add_accent_bar(slide, top=Inches(1.1))

    features = [
        ("Real-Time Monitoring", "USGS + NCS feeds polled every 60s",
         "Earthquake lat/lon/mag/depth from feed", "Live event list + alert banners for M≥6.0",
         "Async polling → normalize → dedup → filter → broadcast"),
        ("Manual Simulation", "User-placed epicenter via POST /api/simulate-earthquake",
         "Lat, Lon, Magnitude, Depth, optional GMPE params", "Grid GeoJSON, contour GeoJSON, district/state summary",
         "Region detection → GMPE select → PGA compute → soil amplify → contour → persist"),
        ("Soil Amplification", "Vs30 GeoTIFF raster → NEHRP site class → amplification factor",
         "Grid cell (lon, lat) → raster sample", "Site class (A-E), soil factor (0.8–1.7), amplified PGA",
         "Raster sample → classify → lookup amplification table → multiply with base PGA"),
        ("Data Export", "GET /api/export/{sim_id}?format=json|csv|geojson",
         "Simulation ID + format", "Full JSON record, CSV district table, or GeoJSON grid",
         "Query DB → serialize to requested format → return as downloadable response"),
    ]

    headers_f = ["Feature", "Purpose", "Inputs", "Outputs", "Internal Processing"]
    add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5), headers_f, features)

    add_speaker_notes(slide,
        "Each major feature follows a clear input-process-output pattern. "
        "Real-time monitoring continuously polls USGS and NCS every 60 seconds. "
        "Manual simulations allow what-if analysis with user-specified parameters. "
        "Soil amplification adjusts base PGA using real Vs30 raster data from national surveys. "
        "Data export supports three formats for downstream analysis.")

    # ================================================================
    # SLIDE 10: Interactive Map Features
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 10)
    add_title_text(slide, "Interactive Hazard Map Features")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(2.4),
        "Map Capabilities",
        ["• WebGL-accelerated vector tile rendering",
         "• Dark and light base map styles (CARTO)",
         "• Satellite and terrain overlays",
         "• India boundary + state boundaries (GeoJSON)",
         "• Smooth zoom/pan with MapLibre GL JS"],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(2.4),
        "Visualization Layers",
        ["• PGA heatmap grid (6,500+ hexagonal cells)",
         "• Contour fill polygons (7 intensity bands)",
         "• Soil amplification overlay (Vs30-based)",
         "• Epicenter marker with pulse animation",
         "• Custom GeoTIFF/raster upload support"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(2.4),
        "User Interactions",
        ["• Click-to-place epicenter on map",
         "• Hover tooltips showing cell-level PGA data",
         "• Layer toggle checkboxes with opacity sliders",
         "• State click → district-level breakdown panel",
         "• Real-time auto-update on live earthquake events"],
        accent_color=ACCENT_ORANGE)

    add_card(slide, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8),
        "PGA Intensity Color Scale (USGS ShakeMap Standard)",
        ["Blue (#3b82f6)     → No Effect    (< 0.02g)",
         "Green (#22c55e)    → Light        (0.02 – 0.115g)",
         "Yellow (#eab308)   → Moderate     (0.115 – 0.215g)",
         "Orange (#f97316)   → Strong       (0.215 – 0.401g)",
         "Dk Orange (#ea580c)→ Very Strong  (0.401 – 0.747g)",
         "Red (#ef4444)      → Severe       (0.747 – 1.39g)",
         "Purple (#7e22ce)   → Violent      (> 1.39g)"],
        accent_color=ACCENT_RED)

    add_card(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.8),
        "Geographic Data & Coordinate System",
        ["• Coordinate Reference System: EPSG:4326 (WGS 84)",
         "• Grid: 20km nationwide hexagonal cells (GeoJSON)",
         "• India boundary: Real GeoJSON polygon (not bounding box)",
         "• Tectonic plates: GeoJSON from Hugo Ahlenius dataset",
         "• Vs30 raster: National GeoTIFF (india_vs30.tif)",
         "• Contour masking: points outside India polygon → NaN"],
        accent_color=ACCENT_CYAN)

    add_speaker_notes(slide,
        "The map interface uses MapLibre GL JS for WebGL-accelerated rendering. "
        "The PGA intensity scale follows the USGS ShakeMap standard with 7 bands from No Effect to Violent. "
        "All geographic data uses EPSG:4326 (WGS 84) coordinates. The contour generation masks out "
        "cells outside India's actual boundary polygon to prevent visual artifacts over oceans and neighboring countries.")

    # ================================================================
    # SLIDE 11: Earthquake Simulation Module
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 11)
    add_title_text(slide, "Earthquake Simulation Module")
    add_accent_bar(slide, top=Inches(1.1))

    add_subtitle_text(slide, "SimulationRunner — app/jobs/simulation_worker.py", top=Inches(1.3), font_size=13, color=MED_GRAY)

    # Simulation pipeline steps
    pipeline = [
        ("1", "Region Detection", "Epicenter → TectonicRegion enum (HIMALAYA / NORTHEAST / PENINSULAR)"),
        ("2", "GMPE Selection", "Region → HimalayanGMPE / NortheastGMPE / PeninsularGMPE (or CustomOverrideGMPE)"),
        ("3", "PGA Computation", "For all 6,500 cells: haversine distance → hypocentral distance → ln(PGA) → exp(ln(PGA))"),
        ("4", "Soil Amplification", "For each cell: Vs30 raster sample → NEHRP class → soil factor → PGA_final = PGA_base × factor"),
        ("5", "Contour Generation", "Scattered PGA → griddata interpolation → Gaussian blur → contourf → GeoJSON polygons"),
        ("6", "Impact Aggregation", "District-level max PGA + severe cell count; state-level avg/max PGA + risk category"),
        ("7", "Persistence", "Save full simulation record to PostgreSQL (grid_geojson, district_summary)"),
        ("8", "Broadcast", "Push simulation_complete message + full results to all connected WebSocket clients"),
    ]

    for i, (num, title, desc) in enumerate(pipeline):
        y = Inches(1.8) + i * Inches(0.65)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_CYAN
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_BG
        p.font.bold = True

        # Title
        txBox = slide.shapes.add_textbox(Inches(1.1), y, Inches(2.5), Inches(0.35))
        tf2 = txBox.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        p2.font.bold = True
        p2.font.name = "Calibri"

        # Description
        txBox2 = slide.shapes.add_textbox(Inches(3.6), y, Inches(9), Inches(0.5))
        tf3 = txBox2.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = LIGHT_GRAY
        p3.font.name = "Calibri"

    # Assumptions box
    add_card(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.9),
        "Assumptions & Limitations",
        ["• Point-source model only (no finite fault rupture)  • No directivity effects  • Empirical GMPE coefficients are approximate regional tunings, not peer-reviewed calibrations  • Depth is user-provided, not independently verified"],
        accent_color=ACCENT_ORANGE)

    add_speaker_notes(slide,
        "The simulation module follows an 8-step pipeline. Each step is modular and testable independently. "
        "The critical performance bottleneck is step 5 (contour generation) which uses matplotlib's contourf — "
        "this is why we wrap it in asyncio.to_thread() to avoid blocking the event loop. "
        "Key assumptions include point-source modeling and empirical coefficient tuning.")

    # ================================================================
    # SLIDE 12: Model Mathematics — GMPE
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 12)
    add_title_text(slide, "Model Mathematics — GMPE")
    add_accent_bar(slide, top=Inches(1.1))

    add_subtitle_text(slide, "Ground Motion Prediction Equation — app/layers/pga/gmpe.py → GenericLogPolynomialGMPE.calculate_pga()",
                      top=Inches(1.3), font_size=12, color=MED_GRAY)

    # Main formula
    add_body_text(slide,
        "ln(PGA) = c₁ + c₂·(M − 6.0) + c₃·(M − 6.0)² + c₄·R − C·ln(R)\n\n"
        "PGA = exp(ln(PGA))   [units: g]",
        left=Inches(0.8), top=Inches(1.8), width=Inches(11), height=Inches(1.0),
        font_size=18, color=ACCENT_CYAN)

    # Variable definitions table
    headers_v = ["Variable", "Definition", "Units", "Source"]
    rows_v = [
        ["M", "Moment Magnitude (Mw)", "dimensionless", "Earthquake feed (USGS/NCS)"],
        ["R", "Hypocentral distance = √(d_surface² + depth²)", "km", "Computed via haversine"],
        ["c₁", "Near-source amplitude constant", "ln(g)", "Region-specific tuning"],
        ["c₂", "Magnitude scaling (linear)", "ln(g)/Mw", "Region-specific tuning"],
        ["c₃", "Magnitude scaling (quadratic)", "ln(g)/Mw²", "Set to 0.0 (all regions)"],
        ["c₄", "Geometric spreading (anelastic)", "ln(g)/km", "Region-specific tuning"],
        ["C", "Logarithmic distance decay", "dimensionless", "Region-specific tuning"],
        ["PGA", "Peak Ground Acceleration output", "g (9.81 m/s²)", "Model output"],
    ]
    add_table(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(2.8), headers_v, rows_v)

    # Regional coefficients
    headers_c = ["Region", "c₁", "c₂", "c₃", "c₄", "C", "File / Class"]
    rows_c = [
        ["Himalaya", "1.40", "0.50", "0.0", "-0.004", "1.0", "HimalayanGMPE"],
        ["Northeast", "1.45", "0.55", "0.0", "-0.0045", "1.0", "NortheastGMPE"],
        ["Peninsular", "1.30", "0.45", "0.0", "-0.003", "0.95", "PeninsularGMPE"],
    ]
    add_table(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2), headers_c, rows_c)

    add_speaker_notes(slide,
        "The GMPE uses a standard logarithmic polynomial form common in seismology. "
        "The equation models PGA as a function of magnitude and hypocentral distance, "
        "with region-specific coefficients that capture differences in crustal structure. "
        "The Himalayan region has the highest c1 (near-source amplitude) due to shallow thrust faulting. "
        "The Peninsular region has a lower C value (0.95 vs 1.0) reflecting slower wave attenuation "
        "in the stable continental shield. NOTE: These coefficients are approximate tunings, not "
        "formally calibrated against recorded strong-motion data. "
        "Source file: app/layers/pga/gmpe.py, lines 40-55.")

    # ================================================================
    # SLIDE 13: Model Mathematics — Soil & Distance
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 13)
    add_title_text(slide, "Model Mathematics — Soil Amplification & Distance")
    add_accent_bar(slide, top=Inches(1.1))

    # Haversine formula
    add_body_text(slide,
        "Haversine Distance Formula — app/layers/pga/engine.py → PGAEngine._haversine_vec()\n\n"
        "a = sin²(Δlat/2) + cos(lat₁)·cos(lat₂)·sin²(Δlon/2)\n"
        "d_surface = R_earth · 2 · arctan2(√a, √(1−a))    [R_earth = 6,371 km]\n\n"
        "Hypocentral Distance:  R = max(√(d_surface² + depth²), 1.0)   [clamped ≥ 1.0 km]",
        left=Inches(0.6), top=Inches(1.4), width=Inches(12), height=Inches(2.0),
        font_size=14, color=ACCENT_CYAN)

    # Soil amplification
    add_body_text(slide,
        "Soil Amplification — app/soil/amplification.py\n\n"
        "PGA_final = PGA_base × Soil_Factor\n\n"
        "Where Soil_Factor is determined by the NEHRP Site Class (derived from Vs30 raster data):",
        left=Inches(0.6), top=Inches(3.5), width=Inches(12), height=Inches(1.5),
        font_size=14, color=ACCENT_GREEN)

    # NEHRP table
    headers_n = ["NEHRP Site Class", "Vs30 Range (m/s)", "Soil Description", "Amplification Factor", "Source File"]
    rows_n = [
        ["A", "> 1500", "Hard Rock", "0.8 (de-amplification)", "soil/site_class.py + soil/amplification.py"],
        ["B", "760 – 1500", "Rock (Reference)", "1.0 (neutral)", "soil/site_class.py + soil/amplification.py"],
        ["C", "360 – 760", "Dense Soil / Soft Rock", "1.2", "soil/site_class.py + soil/amplification.py"],
        ["D", "180 – 360", "Stiff Soil", "1.4", "soil/site_class.py + soil/amplification.py"],
        ["E", "< 180", "Soft Clay", "1.7 (highest)", "soil/site_class.py + soil/amplification.py"],
    ]
    add_table(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), headers_n, rows_n)

    add_speaker_notes(slide,
        "Two additional formulas are critical to the pipeline. "
        "The haversine formula computes great-circle surface distance between the epicenter and each grid cell. "
        "This is combined with depth to get hypocentral distance R, which is clamped to ≥1.0 km to avoid log(0). "
        "The soil amplification follows the NEHRP site classification standard. Vs30 values are sampled from "
        "a national GeoTIFF raster (india_vs30.tif), classified into site classes A-E, and mapped to static "
        "amplification factors. The factors range from 0.8 (hard rock de-amplification) to 1.7 (soft clay amplification). "
        "Reference: NEHRP Recommended Seismic Provisions, FEMA P-750.")

    # ================================================================
    # SLIDE 14: Scientific Foundations
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 14)
    add_title_text(slide, "Scientific & Engineering Foundations")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.2),
        "Seismology Concepts",
        ["• Moment Magnitude (Mw) — standard earthquake",
         "  size measure based on seismic moment",
         "• Hypocentral distance — 3D distance from the",
         "  earthquake focus to a surface observation point",
         "• Ground Motion Prediction Equations (GMPE) —",
         "  empirical relations predicting PGA from M and R",
         "• Tectonic regionalization — classifying zones by",
         "  plate boundary type (thrust, subduction, intraplate)",
         "• Attenuation — decay of seismic wave amplitude",
         "  with distance due to geometric spreading and",
         "  anelastic absorption"],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.2),
        "Hazard Modeling Concepts",
        ["• Peak Ground Acceleration (PGA) — maximum",
         "  ground acceleration during shaking, in units of g",
         "• Site effects / soil amplification — local geology",
         "  modifying ground motion (NEHRP classification)",
         "• Vs30 — time-averaged shear-wave velocity in",
         "  the upper 30 meters of soil/rock",
         "• Contour mapping — spatial interpolation of",
         "  scattered PGA values onto a regular grid",
         "• Gaussian smoothing — noise reduction in",
         "  contour outputs while preserving peak values",
         "• Intensity scales — mapping PGA to human-",
         "  perceived shaking (USGS ShakeMap standard)"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(5.2),
        "Risk Assessment Methods",
        ["• District-level aggregation — max PGA across",
         "  all grid cells within each district",
         "• Severe cell counting — cells with normalized",
         "  PGA > 0.6 flagged as severely affected",
         "• Risk categorization — max PGA mapped to",
         "  categories: LOW / MODERATE / HIGH / SEVERE",
         "  / EXTREME using configurable thresholds",
         "• Damage scoring — composite of average PGA",
         "  (50%) and max PGA (50%), capped at 100",
         "  Formula: min(avg_pga×50 + max_pga×50, 100)",
         "  Source: app/services/impact_aggregator.py L84",
         "• Population exposure — cells with PGA > 0.02g",
         "  contribute their population to affected count"],
        accent_color=ACCENT_ORANGE)

    add_speaker_notes(slide,
        "The application draws on three pillars of earthquake science: "
        "fundamental seismology (magnitude scales, wave propagation, tectonic regimes), "
        "hazard modeling (GMPE, site effects, contour mapping), and risk assessment "
        "(district-level aggregation, severity classification, damage scoring). "
        "The damage score formula is an empirical composite: min(avg_pga*50 + max_pga*50, 100), "
        "found in impact_aggregator.py at line 84.")

    # ================================================================
    # SLIDE 15: Data Pipeline
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 15)
    add_title_text(slide, "Data Pipeline")
    add_accent_bar(slide, top=Inches(1.1))

    pipeline_steps = [
        ("DATA ACQUISITION", "USGS GeoJSON API polled every 60s via aiohttp\nNCS India website scraped for data-json attributes via regex",
         ACCENT_CYAN, Inches(1.6)),
        ("DATA CLEANING", "Validate: -90≤lat≤90, -180≤lon≤180, 0<depth≤700, 1≤mag≤10, age≤2h\nHandle missing depth (default: 10km), mag_type normalization",
         ACCENT_GREEN, Inches(2.8)),
        ("DATA TRANSFORMATION", "normalize_usgs_feature() → canonical EarthquakeEvent dataclass\n_parse_ncs_object() → same canonical format (IST→UTC conversion)",
         ACCENT_ORANGE, Inches(4.0)),
        ("DATA STORAGE", "earthquake_events table (PostgreSQL): source_id, fingerprint, lat, lon, depth, mag\ndedup_cache table: source_id and fingerprint keys for deduplication",
         ACCENT_RED, Inches(5.2)),
        ("DATA SERVING", "REST: GET /api/events (50 most recent), GET /api/events/{id}\nWebSocket: simulation_complete pushed with full GeoJSON payload\nExport: /api/export/{sim_id}?format=json|csv|geojson",
         ACCENT_CYAN, Inches(6.4)),
    ]

    for label, desc, color, top in pipeline_steps:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(2.5), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_BG
        p.font.bold = True

        txBox = slide.shapes.add_textbox(Inches(3.3), top + Inches(0.05), Inches(9.3), Inches(0.85))
        tf2 = txBox.text_frame
        tf2.word_wrap = True
        lines = desc.split("\n")
        for li, line in enumerate(lines):
            if li == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = line
            p2.font.size = Pt(10)
            p2.font.color.rgb = LIGHT_GRAY
            p2.font.name = "Calibri"

    add_speaker_notes(slide,
        "The data pipeline follows a 5-stage process from acquisition to serving. "
        "USGS data arrives as well-structured GeoJSON; NCS data is scraped from HTML data-json attributes. "
        "Both are normalized into the same EarthquakeEvent dataclass. "
        "Deduplication uses two levels: exact source_id match and spatial-temporal fingerprinting "
        "(2-minute time buckets at ~10km lat/lon precision) to catch duplicates across agencies. "
        "The filter stage only passes events M≥4.0 within India's buffered boundary for simulation.")

    # ================================================================
    # SLIDE 16: Data Sources
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 16)
    add_title_text(slide, "Data Sources")
    add_accent_bar(slide, top=Inches(1.1))

    headers_ds = ["Dataset", "Provider", "URL / Source", "Update Freq.", "Purpose"]
    rows_ds = [
        ["2.5+ Hour Feed", "USGS", "earthquake.usgs.gov/.../2.5_hour.geojson", "60s", "Global M≥2.5 earthquakes"],
        ["4.5+ Hour Feed", "USGS", "earthquake.usgs.gov/.../4.5_hour.geojson", "60s", "Fallback feed (larger events)"],
        ["NCS India Events", "NCS/MoES", "riseq.seismo.gov.in/", "60s", "India-specific earthquake feed"],
        ["India Vs30 Raster", "USGS/GSI", "india_vs30.tif (bundled GeoTIFF)", "Static", "Soil shear-wave velocity"],
        ["India Boundary", "Natural Earth", "india_boundary.geojson (bundled)", "Static", "Epicenter validation & masking"],
        ["Tectonic Plates", "Hugo Ahlenius", "TectonicPlateBoundaries.geojson", "Static", "Region classification"],
        ["Nationwide Grid", "Custom", "nationwide_20km.geojson (6,500 cells)", "Static", "Spatial computation grid"],
    ]
    add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(3.0), headers_ds, rows_ds)

    add_card(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(2.5),
        "Data Fields Used Per Source",
        ["USGS GeoJSON: feature.id, geometry.coordinates[lon,lat,depth], properties.mag, properties.magType, properties.place, properties.time, properties.status, properties.alert",
         "NCS HTML: event_id, lat_long, magnitude_depth, origin_time (IST), event_name, event_type",
         "Vs30 Raster: Single band GeoTIFF, CRS EPSG:4326, pixel value = Vs30 in m/s. NoData < 0 treated as 760 m/s (Site Class B)",
         "Grid GeoJSON: Polygon geometry + properties: centroid_lon, centroid_lat, district, state, population"],
        accent_color=ACCENT_CYAN)

    add_speaker_notes(slide,
        "Seven datasets power the application. Two are live feeds (USGS and NCS) polled every 60 seconds. "
        "Five are static datasets bundled with the application. "
        "The Vs30 raster is a critical scientific dataset — it contains shear-wave velocity measurements "
        "for the top 30 meters of soil across India, used to classify soil type and compute amplification factors. "
        "The nationwide grid contains 6,500+ pre-computed hexagonal cells with district and state metadata, "
        "enabling rapid spatial aggregation without runtime geocoding.")

    # ================================================================
    # SLIDE 17: External APIs & Services
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 17)
    add_title_text(slide, "External APIs & Services")
    add_accent_bar(slide, top=Inches(1.1))

    headers_api = ["API / Service", "Endpoint", "Data Retrieved", "Auth", "Rate Limit", "Docs"]
    rows_api = [
        ["USGS Earthquake API", "earthquake.usgs.gov/…/2.5_hour.geojson", "GeoJSON FeatureCollection (M≥2.5)", "None (public)", "~1 req/min", "earthquake.usgs.gov/fdsnws/"],
        ["USGS 4.5+ Feed", "earthquake.usgs.gov/…/4.5_hour.geojson", "GeoJSON FeatureCollection (M≥4.5)", "None (public)", "~1 req/min", "earthquake.usgs.gov/fdsnws/"],
        ["NCS India Website", "riseq.seismo.gov.in/", "HTML with data-json attributes", "None (public)", "~1 req/min", "No formal API docs"],
        ["CARTO Basemap Tiles", "basemaps.cartocdn.com/…", "Vector tiles (dark/light)", "None (public)", "Unlimited (CDN)", "carto.com/basemaps/"],
    ]
    add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(2.2), headers_api, rows_api)

    add_card(slide, Inches(0.6), Inches(4.0), Inches(5.8), Inches(3.2),
        "USGS Feed Integration Details",
        ["• Polled via aiohttp with 15s timeout per request",
         "• 3 retry attempts with exponential backoff (5s, 10s, 15s)",
         "• Automatic fallback from 2.5_hour to 4.5_hour feed",
         "• Consecutive failure counter tracked in poller_stats",
         "• USGS time is milliseconds since epoch → UTC datetime",
         "• Source file: app/ingest/poller.py"],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(6.8), Inches(4.0), Inches(5.8), Inches(3.2),
        "NCS India Scraping Details",
        ["• Website scraped via regex: data-json='\\s*({.*?})\\s*'",
         "• Returns ~150 events per scrape (historical + recent)",
         "• IST timestamps converted to UTC (−5h30m offset)",
         "• Magnitude and depth parsed from 'M: 2.6 , D: 5km'",
         "• No formal API — relies on HTML structure stability",
         "• Source file: app/ingest/ncs_scraper.py"],
        accent_color=ACCENT_GREEN)

    add_speaker_notes(slide,
        "The application integrates with three external services. USGS provides well-structured GeoJSON feeds "
        "with comprehensive earthquake metadata. NCS India does not offer a formal API, so the application "
        "scrapes the website HTML to extract earthquake data from data-json attributes embedded in the page. "
        "CARTO provides free basemap tiles for the map interface. All external calls use async HTTP "
        "with timeouts and retry logic to handle network failures gracefully.")

    # ================================================================
    # SLIDE 18: Database Design
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 18)
    add_title_text(slide, "Database Design")
    add_accent_bar(slide, top=Inches(1.1))

    add_subtitle_text(slide, "PostgreSQL 16 — Schema defined in app/models/repository.py → init_db()",
                      top=Inches(1.3), font_size=12, color=MED_GRAY)

    # Simulations table
    add_card(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.4),
        "simulations",
        ["id SERIAL PRIMARY KEY",
         "timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP",
         "latitude REAL NOT NULL, longitude REAL NOT NULL",
         "magnitude REAL NOT NULL, depth REAL NOT NULL",
         "affected_districts_json JSONB DEFAULT '[]'",
         "grid_geojson_json JSONB DEFAULT NULL",
         "event_id INTEGER DEFAULT NULL",
         "triggered_by TEXT DEFAULT 'manual'"],
        accent_color=ACCENT_CYAN)

    # Events table
    add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.4),
        "earthquake_events",
        ["id SERIAL PRIMARY KEY",
         "source_id TEXT NOT NULL UNIQUE",
         "source TEXT NOT NULL, fingerprint TEXT NOT NULL UNIQUE",
         "latitude REAL, longitude REAL, depth_km REAL",
         "magnitude REAL, mag_type TEXT DEFAULT 'Mw'",
         "origin_time TIMESTAMP NOT NULL",
         "place TEXT, status TEXT, alert_level TEXT",
         "sim_triggered SMALLINT DEFAULT 0"],
        accent_color=ACCENT_GREEN)

    # Dedup table
    add_card(slide, Inches(0.6), Inches(4.5), Inches(5.8), Inches(1.4),
        "dedup_cache",
        ["key TEXT PRIMARY KEY",
         "seen_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP",
         "Used for two-level deduplication: source_id + fingerprint"],
        accent_color=ACCENT_ORANGE)

    # Indexes & Storage
    add_card(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(1.4),
        "Indexes",
        ["idx_simulations_timestamp ON simulations(timestamp DESC)",
         "idx_events_origin_time ON earthquake_events(origin_time DESC)",
         "idx_events_magnitude ON earthquake_events(magnitude DESC)"],
        accent_color=ACCENT_RED)

    add_card(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.9),
        "Storage Strategy",
        ["Connection pool: ThreadedConnectionPool(minconn=2, maxconn=10). Context manager (_get_conn) with auto-commit/rollback.  JSONB for grid data enables future spatial queries. Daily cleanup task removes data older than 7 days."],
        accent_color=MED_GRAY)

    add_speaker_notes(slide,
        "The database schema uses three tables. The simulations table stores full simulation results "
        "including the complete grid GeoJSON as JSONB, enabling future spatial queries. "
        "The earthquake_events table tracks all ingested events with unique constraints on source_id "
        "and fingerprint for deduplication. The dedup_cache provides fast lookup for the two-level "
        "deduplication strategy. A ThreadedConnectionPool with 2-10 connections handles concurrent access "
        "from the FastAPI thread pool. A daily cleanup task runs every 24 hours to remove old data.")

    # ================================================================
    # SLIDE 19: Core Algorithms
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 19)
    add_title_text(slide, "Core Algorithms")
    add_accent_bar(slide, top=Inches(1.1))

    algos = [
        ("Vectorized PGA Calculation", "PGAEngine.compute()", "app/layers/pga/engine.py",
         "Extracts all cell coordinates as NumPy arrays, computes haversine distances vectorized, "
         "calculates hypocentral distance R, delegates to GMPE.calculate_pga() — all without Python for-loops."),
        ("Spatial-Temporal Deduplication", "compute_fingerprint()", "app/ingest/normalizer.py",
         "SHA-256 hash of (lat rounded to 0.1°, lon rounded to 0.1°, mag rounded to 0.1, 2-minute time bucket). "
         "Catches the same earthquake reported by different agencies with slight parameter variations."),
        ("Contour Generation", "generate_contour_geojson()", "app/services/contour_generator.py",
         "1) scipy.griddata() for scattered-to-regular interpolation 2) Gaussian filter (σ=1.5) for smoothing "
         "3) Peak amplitude restoration post-blur 4) India polygon masking 5) matplotlib contourf → GeoJSON via geojsoncontour."),
        ("Tectonic Region Detection", "get_tectonic_region()", "app/layers/pga/regions.py",
         "Loads real plate boundary GeoJSON, buffers by 4.0° (~440km), checks point-in-polygon. "
         "Longitude > 89° splits Himalaya from Northeast. Latitude < 25° overrides to Peninsular."),
        ("Risk Classification", "_classify_risk()", "app/services/impact_aggregator.py",
         "Threshold-based: max_pga > 0.8 → EXTREME, > 0.6 → SEVERE, > 0.4 → HIGH, > 0.2 → MODERATE, else LOW. "
         "Damage score = min(avg_pga×50 + max_pga×50, 100)."),
    ]

    headers_a = ["Algorithm", "Function", "Source File", "Description"]
    add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5), headers_a, algos)

    add_speaker_notes(slide,
        "Five core algorithms power the application. The vectorized PGA calculation avoids Python for-loops "
        "entirely, achieving ~40-80x speedup over the naive implementation. The deduplication algorithm uses "
        "a spatial-temporal fingerprint with 2-minute time buckets and ~10km spatial resolution. "
        "The contour generation pipeline is the most complex, combining scattered interpolation, Gaussian smoothing, "
        "peak restoration, geographic masking, and format conversion in a single pass.")

    # ================================================================
    # SLIDE 20: Map Rendering Process
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 20)
    add_title_text(slide, "Map Rendering Process")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(2.4),
        "Geographic Data Flow",
        ["1. Backend computes PGA for all grid cells",
         "2. Contour polygons generated as GeoJSON",
         "3. Full grid + contour GeoJSON sent via WS",
         "4. Frontend adds as MapLibre GL sources",
         "5. fill-extrusion / fill layers render on GPU"],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(2.4),
        "Tile & Layer Management",
        ["• Base tiles: CARTO dark-matter / positron",
         "• Overlay layers managed via mapLayerService.js",
         "• Layer ordering: base → boundary → grid → contour",
         "• Opacity controlled per-layer via Zustand state",
         "• Custom raster upload via rasterService.js"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(2.4),
        "Coordinate Systems",
        ["• All data in EPSG:4326 (WGS 84)",
         "• MapLibre internally uses Web Mercator",
         "  (EPSG:3857) for tile rendering",
         "• Contour coordinates: lon/lat in degrees",
         "• Raster bounds: rasterio BoundingBox (lon/lat)"],
        accent_color=ACCENT_ORANGE)

    add_card(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.8),
        "Contour Rendering Pipeline Detail",
        ["Step 1: Extract (lon, lat, PGA) triples from 6,500 grid cells — scattered points",
         "Step 2: Build 400×400 regular meshgrid over India's bounding box using np.mgrid",
         "Step 3: Interpolate PGA onto meshgrid via scipy.interpolate.griddata(method='linear')",
         "Step 4: Apply Gaussian blur (σ=1.5) to smooth jagged edges, then restore peak amplitude",
         "Step 5: Build GeoPandas mask — set NaN for meshgrid points outside India boundary polygon",
         "Step 6: Run matplotlib.contourf() with 7 PGA levels + dynamic upper bound",
         "Step 7: Convert matplotlib contour to GeoJSON via geojsoncontour.contourf_to_geojson()",
         "Step 8: Frontend receives GeoJSON → adds as MapLibre GL fill layer with per-band colors"],
        accent_color=ACCENT_RED)

    add_speaker_notes(slide,
        "Map rendering involves a coordinated pipeline between backend and frontend. "
        "The backend generates contour polygons as GeoJSON using an 8-step process involving "
        "scattered interpolation, Gaussian smoothing, geographic masking, and matplotlib contourf. "
        "The frontend receives this via WebSocket and renders it as a MapLibre GL fill layer. "
        "The 400x400 meshgrid resolution provides a good balance between visual quality and computation time. "
        "The India boundary mask prevents contour artifacts over the ocean and neighboring countries.")

    # ================================================================
    # SLIDE 21: User Workflow Example
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 21)
    add_title_text(slide, "User Workflow — Complete Example")
    add_accent_bar(slide, top=Inches(1.1))

    add_subtitle_text(slide, "Scenario: Manual simulation of M6.5 earthquake near Shimla (31.1°N, 77.1°E, depth 10km)",
                      top=Inches(1.3), font_size=13, color=MED_GRAY)

    workflow = [
        ("User clicks map at (31.1°N, 77.1°E)", "Browser sends lat/lon → backend GET /api/region → returns 'HIMALAYA'"),
        ("User sets M=6.5, depth=10km, clicks Run", "POST /api/simulate-earthquake {magnitude:6.5, depth:10, lat:31.1, lon:77.1}"),
        ("Backend validates epicenter", "is_epicenter_valid() checks Point(77.1, 31.1) inside BUFFERED_INDIA (9° buffer)"),
        ("GMPESelector picks HimalayanGMPE", "c1=1.40, c2=0.50, c3=0.0, c4=-0.004, C=1.0"),
        ("PGAEngine computes PGA for 6,500 cells", "Haversine distances → R_hypo → ln(PGA) = 1.40+0.5(M-6)-0.004R-ln(R) → exp()"),
        ("SoilEngine enriches each cell", "Vs30 sampled from india_vs30.tif → site class → amplification factor"),
        ("PGA_final = PGA_base × soil_factor", "Example: cell at 31.0°N has Vs30=280 → Class D → factor 1.4 → PGA 0.35→0.49g"),
        ("ContourGenerator produces GeoJSON", "griddata interpolation → Gaussian blur → India mask → contourf → 7-band polygons"),
        ("ImpactAggregator builds summaries", "Districts: Shimla max_pga=0.49g EXTREME | States: Himachal Pradesh risk=EXTREME"),
        ("Result saved to PostgreSQL + returned", "sim_id=42, grid_geojson (6,500 features), contour_geojson, district/state summary"),
    ]

    for i, (action, detail) in enumerate(workflow):
        y = Inches(1.7) + i * Inches(0.55)
        # Step number
        txN = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.4), Inches(0.4))
        tf = txN.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}."
        p.font.size = Pt(10)
        p.font.color.rgb = ACCENT_CYAN
        p.font.bold = True
        p.font.name = "Calibri"

        # Action
        txA = slide.shapes.add_textbox(Inches(0.9), y, Inches(4.5), Inches(0.5))
        tf2 = txA.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = action
        p2.font.size = Pt(10)
        p2.font.color.rgb = WHITE
        p2.font.bold = True
        p2.font.name = "Calibri"

        # Detail
        txD = slide.shapes.add_textbox(Inches(5.5), y, Inches(7.3), Inches(0.5))
        tf3 = txD.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = detail
        p3.font.size = Pt(9)
        p3.font.color.rgb = LIGHT_GRAY
        p3.font.name = "Calibri"

    add_speaker_notes(slide,
        "This slide traces a complete manual simulation workflow with concrete values. "
        "A user places an epicenter near Shimla and runs an M6.5 simulation. "
        "The system detects the Himalayan tectonic region, selects the HimalayanGMPE, "
        "computes PGA for all 6,500 grid cells, applies Vs30 soil amplification, "
        "generates contour polygons, aggregates district and state impacts, "
        "saves to the database, and returns the full result — all in under 10 seconds. "
        "The example shows a cell with Vs30=280 m/s (Class D, stiff soil) amplifying PGA from 0.35g to 0.49g.")

    # ================================================================
    # SLIDE 22: Performance & Scalability
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 22)
    add_title_text(slide, "Performance & Scalability")
    add_accent_bar(slide, top=Inches(1.1))

    headers_p = ["Component", "Technique", "Impact"]
    rows_p = [
        ["PGA Computation", "Vectorized NumPy (no Python loops)", "~40-80x speedup over per-cell loop"],
        ["Vs30 Sampling", "Batch rasterio.sample() per state group", "1 I/O call instead of 6,500"],
        ["Grid Loading", "Pre-loaded once at startup, deep-copied per sim", "Zero disk I/O per simulation"],
        ["Contour Generation", "asyncio.to_thread() for blocking matplotlib", "Event loop stays free for WS/poller"],
        ["Contour Masking", "India boundary mask cached after first call", "400×400 mask computed only once"],
        ["DB Connections", "ThreadedConnectionPool(min=2, max=10)", "No connection overhead per query"],
        ["Raster Handles", "Opened once at startup, kept in memory", "No rasterio.open() in hot path"],
        ["Deduplication", "PostgreSQL-backed (no Redis dependency)", "Persistent across restarts"],
        ["Poller Network", "aiohttp with timeout + exponential backoff", "Graceful degradation on failure"],
        ["WebSocket", "Set-based client tracking, dead conn removal", "O(1) add/remove, batch broadcast"],
    ]
    add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(4.0), headers_p, rows_p)

    add_card(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(1.3),
        "Response Time Estimates",
        ["• PGA computation (6,500 cells): ~50-100ms  •  Soil amplification: ~20-50ms  •  Contour generation: ~2-3s (blocking, threaded)",
         "• Impact aggregation: ~10ms  •  DB persistence: ~5ms  •  WebSocket broadcast: ~1ms  •  Total end-to-end: ~3-5 seconds per simulation"],
        accent_color=ACCENT_GREEN)

    add_speaker_notes(slide,
        "Performance optimization is critical because the system needs to generate hazard maps in near-real-time. "
        "The biggest win is vectorized NumPy computation — processing 6,500 grid cells takes ~50-100ms vs ~5 seconds "
        "with Python for-loops. The contour generation step (matplotlib contourf) is the bottleneck at 2-3 seconds, "
        "which is why it's wrapped in asyncio.to_thread() to avoid blocking the event loop. "
        "Caching strategies (grid pre-loading, raster handle pooling, contour mask caching) eliminate redundant I/O.")

    # ================================================================
    # SLIDE 23: Security Considerations
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 23)
    add_title_text(slide, "Security Considerations")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.5),
        "Input Validation",
        ["• Pydantic model validation on all API inputs:",
         "  - magnitude: ge=1.0, le=10.0",
         "  - depth: gt=0, le=10000",
         "  - latitude: ge=-90.0, le=90.0",
         "  - longitude: ge=-180.0, le=180.0",
         "• Epicenter validated against India boundary polygon",
         "• EarthquakeEvent fields validated in is_valid()"],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5),
        "API Protection",
        ["• WebSocket DoS guard: MAX_CLIENTS = 100",
         "• Async queue maxsize = 500 (prevents memory exhaustion)",
         "• CORS configured (currently allow_origins=['*'])",
         "• Dead WebSocket connections auto-cleaned",
         "• Poller timeout: 15s per request, 3 retries max",
         "• Queue overflow: events dropped with warning log"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(0.6), Inches(4.4), Inches(5.8), Inches(2.5),
        "Database Security",
        ["• Parameterized SQL queries (%s placeholders) —",
         "  prevents SQL injection",
         "• Connection pool with auto-rollback on exceptions",
         "• Database credentials via environment variable",
         "  (DATABASE_URL in docker-compose.yml)",
         "• No raw SQL string concatenation anywhere"],
        accent_color=ACCENT_ORANGE)

    add_card(slide, Inches(6.8), Inches(4.4), Inches(5.8), Inches(2.5),
        "Areas for Improvement",
        ["• No authentication/authorization currently —",
         "  all endpoints are public",
         "• CORS is fully open (allow_origins=['*']) —",
         "  should be restricted in production",
         "• No rate limiting on simulation endpoint",
         "• Database password is hardcoded in compose file",
         "• No HTTPS enforcement (relies on reverse proxy)"],
        accent_color=ACCENT_RED)

    add_speaker_notes(slide,
        "Security is implemented at multiple layers but has room for improvement. "
        "Input validation uses Pydantic models with strict bounds checking. "
        "SQL injection is prevented by parameterized queries throughout. "
        "DoS protection exists for WebSocket connections and the simulation queue. "
        "However, there is currently no authentication, rate limiting, or HTTPS enforcement — "
        "these would need to be added for a production deployment. "
        "CORS should be restricted to known frontend origins.")

    # ================================================================
    # SLIDE 24: Challenges Encountered
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 24)
    add_title_text(slide, "Challenges Encountered")
    add_accent_bar(slide, top=Inches(1.1))

    challenges = [
        ("Matplotlib Blocking Event Loop", "Technical",
         "matplotlib's contourf is synchronous and takes 2-3s — blocked WebSocket pings and poller",
         "Wrapped in asyncio.to_thread() + used Agg backend for thread safety"),
        ("Cross-Source Deduplication", "Data",
         "Same earthquake reported by USGS and NCS with slightly different parameters",
         "Spatial-temporal fingerprinting: SHA-256 of (lat±0.1°, lon±0.1°, mag±0.1, 2-min bucket)"),
        ("NCS India No API", "Data",
         "NCS does not provide a formal REST API for earthquake data",
         "Built HTML scraper using regex to extract data-json attributes from the website"),
        ("WebSocket Variable Scoping", "Technical",
         "Python's _CLIENTS -= dead creates a local variable reference, causing UnboundLocalError",
         "Changed to _CLIENTS.difference_update(dead) for in-place set mutation"),
        ("Tectonic Region Boundaries", "Scientific",
         "Simple lat/lon bounding boxes misclassified regions near zone boundaries",
         "Used real tectonic plate boundary GeoJSON with 4° buffer + longitude/latitude overrides"),
        ("Contour Artifacts", "Performance",
         "Gaussian blur reduced peak PGA values in contour maps",
         "Added peak amplitude restoration: grid_z *= original_max / new_max after blur"),
    ]

    headers_ch = ["Challenge", "Type", "Problem", "Solution"]
    add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5), headers_ch, challenges)

    add_speaker_notes(slide,
        "Six significant challenges were encountered during development. "
        "The matplotlib event loop blocking issue was the most impactful — it caused WebSocket pings to drop "
        "and the poller to stall during contour generation. The solution (asyncio.to_thread) is now a documented "
        "best practice in the codebase. The NCS scraping challenge required building a fragile HTML parser "
        "that depends on the website's DOM structure remaining stable. The WebSocket scoping bug (_CLIENTS -= dead) "
        "was a subtle Python gotcha that caused all broadcasts to fail silently.")

    # ================================================================
    # SLIDE 25: Limitations
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 25)
    add_title_text(slide, "Current Limitations")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.0),
        "Scientific Limitations",
        ["• GMPE coefficients are approximate tunings,",
         "  not peer-reviewed calibrations against",
         "  recorded strong-motion data",
         "• Point-source earthquake model only —",
         "  no finite fault rupture modeling",
         "• No directivity effects (azimuthal variation)",
         "• Static soil amplification factors (no PGA-",
         "  dependent nonlinear site response)",
         "• No basin depth (Z1.0, Z2.5) correction",
         "• c₃ = 0 for all regions (no quadratic",
         "  magnitude saturation term)"],
        accent_color=ACCENT_ORANGE)

    add_card(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.0),
        "Technical Limitations",
        ["• Single-worker simulation queue (no parallel",
         "  processing of multiple earthquakes)",
         "• NCS scraper depends on HTML structure",
         "  (fragile, may break if website changes)",
         "• No authentication or rate limiting",
         "• Contour resolution fixed at 400×400 grid",
         "• No caching of simulation results on frontend",
         "• PostgreSQL stores full GeoJSON as JSONB",
         "  (large storage footprint per simulation)"],
        accent_color=ACCENT_RED)

    add_card(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(5.0),
        "Known Issues",
        ["• CORS is fully open (allow_origins=['*'])",
         "• Database credentials hardcoded in compose",
         "• No graceful handling of raster file corruption",
         "• Poller failure counter not exposed to frontend",
         "• No unit tests for GMPE calculations",
         "• No automated integration testing",
         "• Grid cell population data not verified",
         "• IST→UTC conversion assumes fixed offset",
         "  (no DST awareness, though India doesn't",
         "  observe DST)"],
        accent_color=MED_GRAY)

    add_speaker_notes(slide,
        "The most significant limitation is the GMPE calibration — the coefficients are approximate tunings "
        "rather than formal calibrations against recorded strong-motion data from Indian earthquakes. "
        "An earlier attempt to implement the BSSA14 (Boore et al. 2014) peer-reviewed GMPE was explored "
        "but not completed. The single-worker queue means only one simulation runs at a time — "
        "during earthquake swarms, events would queue up. The NCS scraper is inherently fragile.")

    # ================================================================
    # SLIDE 26: Key Learnings
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 26)
    add_title_text(slide, "Key Learnings")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.2),
        "Engineering Lessons",
        ["• asyncio + to_thread() is essential when mixing",
         "  async I/O with blocking scientific libraries",
         "• NumPy vectorization provides dramatic speedups",
         "  — always profile before using Python loops",
         "• WebSocket state management requires careful",
         "  Python variable scoping (global vs. local)",
         "• Docker Compose with volume mounts enables",
         "  hot-reload without container rebuilds",
         "• Module-level imports of heavy libraries",
         "  (matplotlib, rasterio) minimize cold starts",
         "• Connection pooling is critical for threaded",
         "  web servers — never create connections per-request"],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.2),
        "Domain Knowledge",
        ["• Earthquake magnitude is logarithmic — each",
         "  unit represents ~31.6x more energy released",
         "• Vs30 is the single most important site parameter",
         "  for predicting ground motion amplification",
         "• NEHRP site classification is a widely adopted",
         "  standard but simplifies complex soil behavior",
         "• Tectonic region matters enormously — Himalayan",
         "  earthquakes attenuate differently than Peninsular",
         "• Contour smoothing must preserve peak values",
         "  — naive Gaussian blur underestimates maxima",
         "• Real-time is not just fast computation — it's also",
         "  fast data acquisition + push notification"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(5.2),
        "Development Insights",
        ["• Modular architecture (BaseGMPE, BaseHazardLayer)",
         "  enables swapping models without changing pipeline",
         "• Canonical data models (EarthquakeEvent) simplify",
         "  multi-source ingestion — normalize early, process uniformly",
         "• Logging > print() — structured logging with module",
         "  names enables filtering in production",
         "• Git rollback is essential when experimenting with",
         "  complex scientific implementations (BSSA14 attempt)",
         "• Browser DevTools WebSocket inspector is invaluable",
         "  for debugging real-time data flow",
         "• GeoJSON is the lingua franca of web GIS —",
         "  design every pipeline stage to produce/consume it"],
        accent_color=ACCENT_ORANGE)

    add_speaker_notes(slide,
        "Three categories of learning emerged from this project. Engineering lessons center on the challenges "
        "of combining async web frameworks with blocking scientific libraries. Domain knowledge about seismology "
        "and soil mechanics was essential for building a scientifically credible application. "
        "Development insights include the importance of modular architecture for scientific applications, "
        "canonical data models for multi-source ingestion, and proper logging for debugging distributed systems.")

    # ================================================================
    # SLIDE 27: Conclusion
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 27)
    add_title_text(slide, "Conclusion")
    add_accent_bar(slide, top=Inches(1.1))

    add_card(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(2.0),
        "Summary",
        ["HazardMap is a full-stack real-time earthquake hazard assessment platform covering the Indian subcontinent.",
         "It combines live data ingestion (USGS + NCS India), region-specific ground motion prediction, Vs30-based soil amplification,",
         "smooth contour visualization, and WebSocket-based real-time updates — all in a Docker-deployed web application.",
         "The system processes 6,500+ grid cells per simulation in under 10 seconds end-to-end."],
        accent_color=ACCENT_CYAN)

    add_card(slide, Inches(0.6), Inches(3.9), Inches(5.8), Inches(3.0),
        "Impact",
        ["• Demonstrates feasibility of browser-based",
         "  real-time seismic hazard assessment",
         "• Provides a framework for integrating multiple",
         "  earthquake data sources with deduplication",
         "• Shows how scientific computing (NumPy, SciPy,",
         "  matplotlib) can be effectively integrated into",
         "  modern async web architectures",
         "• Enables rapid what-if analysis for disaster",
         "  preparedness and emergency response planning"],
        accent_color=ACCENT_GREEN)

    add_card(slide, Inches(6.8), Inches(3.9), Inches(5.8), Inches(3.0),
        "Future Potential",
        ["• Implement BSSA14 peer-reviewed GMPE with",
         "  validated OpenQuake coefficients",
         "• Add real-time ShakeMap integration from USGS",
         "• Machine learning-based rapid damage estimation",
         "• Multi-hazard layers: tsunami, landslide, liquefaction",
         "• Mobile app with push notifications for alerts",
         "• Historical earthquake database + replay mode",
         "• Integration with building vulnerability databases"],
        accent_color=ACCENT_ORANGE)

    add_speaker_notes(slide,
        "In conclusion, HazardMap demonstrates that a scientifically grounded, real-time earthquake hazard "
        "assessment platform can be built using modern open-source technologies. The modular architecture "
        "is designed for extensibility — future work could include peer-reviewed GMPE implementations, "
        "multi-hazard modeling, and machine learning integration. The most impactful near-term improvement "
        "would be implementing the BSSA14 GMPE with validated coefficients from OpenQuake.")

    # ================================================================
    # SLIDE 28: References
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_section_number(slide, 28)
    add_title_text(slide, "References")
    add_accent_bar(slide, top=Inches(1.1))

    refs = (
        "Scientific References:\n"
        "• NEHRP Recommended Seismic Provisions, FEMA P-750 (2020) — Site classification standard\n"
        "• Boore, D.M., Stewart, J.P., Seyhan, E., Atkinson, G.M. (2014) — BSSA14 GMPE (NGA-West2)\n"
        "• USGS ShakeMap Manual — PGA intensity color scale and thresholds\n"
        "• Haversine Formula — Great-circle distance computation (standard geodesy)\n\n"
        "Datasets:\n"
        "• USGS Earthquake Hazards Program — earthquake.usgs.gov/earthquakes/feed/\n"
        "• National Center for Seismology, MoES India — riseq.seismo.gov.in/\n"
        "• USGS Global Vs30 Model — earthquake.usgs.gov/data/vs30/\n"
        "• Natural Earth — India administrative boundaries\n"
        "• Hugo Ahlenius / Nordpil — Tectonic plate boundaries GeoJSON\n\n"
        "Libraries & Frameworks:\n"
        "• FastAPI (tiangolo) — fastapi.tiangolo.com\n"
        "• React 19 — react.dev\n"
        "• MapLibre GL JS — maplibre.org\n"
        "• NumPy / SciPy / matplotlib — numpy.org / scipy.org / matplotlib.org\n"
        "• Shapely / GeoPandas / rasterio — shapely.readthedocs.io / geopandas.org / rasterio.readthedocs.io\n"
        "• PostgreSQL 16 — postgresql.org\n"
        "• geojsoncontour — github.com/bartromgens/geojsoncontour\n"
        "• Zustand — github.com/pmndrs/zustand"
    )

    add_body_text(slide, refs, left=Inches(0.6), top=Inches(1.5), width=Inches(12), height=Inches(5.5),
                  font_size=11, color=LIGHT_GRAY)

    add_speaker_notes(slide,
        "This slide lists all references used in the application, organized into three categories: "
        "scientific references for the seismological methods, datasets for the input data, "
        "and libraries and frameworks for the implementation. "
        "Key scientific references include the NEHRP site classification standard and the USGS ShakeMap manual. "
        "The BSSA14 paper is listed as a reference for future GMPE implementation work.")

    # ── Save ──────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__), "HazardMap_Presentation.pptx")
    prs.save(output_path)
    print(f"\n✅ Presentation saved to: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    build_presentation()
